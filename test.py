import http.server
import socketserver
import pickle
import os
import cgi
import pypdf
import docx
from pptx import Presentation
from urllib.parse import urlparse, parse_qs, quote, unquote
from pathlib import Path
port = 2704
Handler = http.server.SimpleHTTPRequestHandler

from http.server import HTTPServer, BaseHTTPRequestHandler

# Define a simple class to hold vocab items
class Root:
    def getVerbForms(self, rt: str, can_have_objects):
        verb_forms = set()
        #INFINITIVE 
        verb_forms.add("ku" + rt)
        #HABITUAL
        verb_forms.add("hu" + rt)
        #INDICATIVES and RELATIVES
        # there are only 11 monosyllabic acting verbs. 
        isMonosyllabic = rt in ["la","fa","ja","nywa","la",'pa','wa','cha','chwa','isha','nya']
        #add for looping through; positive verb_forms
        positiveSubjects = ["ni","u","a","tu","wa","m","u","i","li","ya","ki","vi","zi","ku"]
        #covers present, past, future, perfect, perfect already, conditional, regret, and narriative
        positiveTenses = ["na","li","ta","me", "mesha","nge","ngali", "ka"]
        #relative markers, including blank which would be the indicative
        relatives = ['','ye','o','cho','vyo','yo','lo','zo','ko','po','mo']
        #affirmative indicatives, and relatives as only the affirmative can have relatives. 
        affirmativeRoot = "ku" + rt if isMonosyllabic else rt
        for i in positiveSubjects:
            for j in positiveTenses:
                for k in relatives:
                    verb_forms.add( i + j + k + affirmativeRoot )
        #temporal relative
        for i in positiveSubjects:
            for j in ['li','taka','na']:
                verb_forms.add(i + j + "po" + rt)
        #add for looping through; negative verb_forms
        negativeSubjects = ['si','hu','ha','hatu','hawa','ham','hau','hai','hali','haya','haki','havi','hazi','haku']
        #negative indicative in the present
        negativeRoot = rt[:-1] + "i" if rt[-1] == "a" else rt
        for i in negativeSubjects:
            verb_forms.add(i + negativeRoot)
        #past, future, perfect, conditional, and regret tenses
        negativeTenses = ["ku","ja","ta","singe","singali"]
        for i in negativeSubjects:
            for j in negativeTenses:
                verb_forms.add(i + j + rt)
        #SUBJUNCTIVES
        subjunctiveRoot = rt[:-1] + "e" if rt[-1] == "a" else rt
        for i in positiveSubjects:
            verb_forms.add(i + subjunctiveRoot)
        #CONDITIONALS
        for i in positiveSubjects:
            verb_forms.add(i + "ki" + rt)
        #IMPERATIVES
        #singular, informal
        verb_forms.add(affirmativeRoot)
        #plural formal
        verb_forms.add(rt[:-1] + "eni" if rt[-1] == "a" else rt + "ni")
        #OBJECTS AND RECIPROCALS, starting with personal pronouns and moving to noun class infixes.
        if can_have_objects:
            objectInfixes = ["ni","ku","m","wa","tu","ki","vi","u","i","li","ya","zi","ku","pa"]
            #If it starts with a vowel,
            if rt[0] in "aeiou":
                objectInfixes[2] = "mw"
            #runs the function again, but with objects turned off and the root replaced with the verb with object
            for i in objectInfixes:
                verb_forms.update(getVerbForms(i + rt, can_have_objects=False))        
            verb_forms.update(getVerbForms(rt + "na" if rt[-1] == "a" else rt + "ana", can_have_objects=False))

        return verb_forms 
    def __init__(self, root, pos, filename, include_in_vocab_list=True, english=None, plural=None, transitive=False):
        self.root = root
        self.pos = pos
        self.filename = filename
        self.english = english
        self.plural = plural
        self.transitive = transitive
        self.include_in_vocab_list = include_in_vocab_list
        self.forms = set()

        if pos == "verb":
            self.forms = self.getVerbForms(root, can_have_objects=transitive)
        
        #nouns are bare simple, with four forms: singular, plural, singular with ni, and plural with ni.
        if pos == "noun":
            self.forms = {root, plural, root + "ni", plural + "ni"}

        if pos == "adjective":
            # n is at the beginning to deal with the two below rules
            prefixes = ["n","m","wa","ki","vi","mi","","ma","ku","pa"]
            # chief pocket rule makes nfupi -> fupi
            chief_pocket = root[0] in "fpkt" or root[:2] in ["ch"]
            if chief_pocket:
                prefixes.pop(0)
            # if it starts with a vowel, the m prefix becomes mw
            if root[0] in "aeio":
                prefixes[1] = "mw"
            # r rule makes nrefu to ndefu
            if root[0] == "r":
                prefixes.pop(0)
                self.forms.add("nd" + root[1:])
            for i in prefixes:
                self.forms.add(i + root)

        if pos == "other":
            self.forms.add(root)
    
    def get_forms(self):
        return self.forms


class Vocab:
    def __init__(self, name):
        self.name = name
        self.roots = []
        self.forms = {'wa','cha','vya','ya','la','kwa','pa','mwa','za','huyu','huyo','yule','hawa','hao','wale','hiki','hicho','kile','hivi','hivyo','vile','huu','huo','ule','hii','hiyo','ile','hili','hilo','lile','haya','hayo','yale','hii','hiyo','ile','hizi','hizo','zile','huu','huo','ule','hizi','hizo','zile','huku','huko','kule','hapa','hapo','pale','humu','humo','mle','ambaye','ambao','ambacho','ambavyo','ambao','ambayo','ambalo','ambayo','ambazo','ambako','ambapo','ambamo','mimi','wewe','yeye','sisi','ninyi','wao','mwenye','wenye','yenye','lenye','chenye','vyenye','zenye','penye','kwenye','mwenye','nina','una','ana','tuna','mna','wana','ina','lina','ana','kina','vina','ina','zina','pana','kuna','sina','huna','hana','hatuna','hamna','hawana','hauna','haina','halina','hayana','hakina','hvina','haina','hazina','hauna','hakuna','hamna','yeyote','wowote','yoyote','lolote','chochote','vyovyote','zozote','popote','kokote','momote','yupi','wepi','upi','ipi','lipi','yapi','kipi','vipi','ipi','zipi','upi','zipi','papi','kupi','mpi'}
    def in_vocab(self, word):
        return word in self.forms
    def add_root(self, root: Root):
        self.roots.append(root)
        self.forms.update(root.get_forms())
    def get_vocab_list(self, filename = None):
        if filename is not None:
            return [[i.root, i.english] for i in self.roots if i.filename == filename and i.include_in_vocab_list]
        else:
            return [[i.root, i.english] for i in self.roots if i.include_in_vocab_list]
    def sanitize_text(self, text):
        #returns text with existing forms removed.
        #text should be preprocessed, with punc removed, all lowercase, and separated by spaces.
        words = text.split()
        sanitized_words = [word for word in words if word not in self.forms]
        return ' '.join(sanitized_words)
    def fast_sanitize(self, text, root):
        #removes only one root from the text
        words = text.split()
        sanitized_words = [word for word in words if word not in root.get_forms()]
        return ' '.join(sanitized_words)



# Functions to read and write the class list to a file
def write_to_file():
    with open("items.pkl", "wb") as f:
        pickle.dump(classList, f)

# Function to read the class list from a file
def read_from_file():
    with open("items.pkl", "rb") as f:
        a = pickle.load(f)
        return a

# Initialize the class list, but grab from file if it exists
try:
    classList = read_from_file()

except FileNotFoundError:
    classList = []
    pickle.dump(classList, open("items.pkl", "wb"))

#Initialize directory for document uploads
os.makedirs("uploads", exist_ok=True)

# Function to generate HTML list items from the class list
def make_html_list():
    classList = read_from_file()
    text = ""
    for i in range(len(classList)):
        name = classList[i].name
        #GREAT MERGER
        text += f""" 
        <li>
            <a href="/class?{i}">{name}
                <form method="post" action="/" onsubmit="return confirm('Are you sure you want to delete this class?');">
                    <input type="hidden" name="action" value="delete">
                    <input type="hidden" name="id" value="{i}">
                    <button type="submit">Delete</button>
                </form>
            </a>
        </li>
        """
    return text

def make_file_html_list(class_name, class_id):
    # Grabs uploaded files
    try:
        file_array = os.listdir(f"uploads{os.sep}class_{class_name}")
    except FileNotFoundError:
        file_array = []

    file_list_html = ""

    # Makes <li> of uploaded files (Delete Form Inside HTML Text)
    for file in file_array:
        safe = quote(file)
        file_list_html += (f"""
                            <li>
                            <a href="/download?class={class_name}&file={safe}">{file}</a>
                            <form method="post" action="/" onsubmit="return confirm('Are you sure you want to delete this file?');">
                                <input type="hidden" name="action" value="delete_file">
                                <input type="hidden" name="classname" value="{class_name}">
                                <input type="hidden" name="id" value="{file_array.index(file)}">
                                <input type="hidden" name="classID" value="{class_id}">

                                <button type="submit">Delete</button>
                            </form>
                            </li>""")
    return file_list_html

class SimpleHandler(BaseHTTPRequestHandler):

    # Handle Get requests, different pages
    def do_GET(self):
        
        page = urlparse(self.path)

        # Serve CSS
        if page.path == "/style.css":
            self.send_response(200)
            self.send_header("Content-type", "text/css")
            self.end_headers()
            with open("style.css", "rb") as f:
                self.wfile.write(f.read())
        # Serve class page
        elif page.path == "/class":
            self.send_response(200)
            self.send_header("Content-type", "text/html")
            self.end_headers()
            # Access class.html filee
            with open("class.html", "rb") as f:
                html = f.read()

            # Sets values for html replacements
            class_number = page.query
            class_name = classList[int(class_number)].name

            # Handles title of class page
            html = html.replace(b"<!-- c -->", class_name.encode('utf-8'))

            # Handles form building for uploads
            html = html.replace(b"replace_with_class_name", class_name.encode('utf-8'))

            file_list_html = make_file_html_list(class_name, class_number)

            # # Grabs uploaded files

            # And inserts the <li>s onto the page
            html = html.replace(b"<!-- FILES -->", file_list_html.encode('utf-8'))
            html = html.replace(b"CLASSIDPLACEHOLDER", class_number.encode('utf-8'))
            self.wfile.write(html)

        elif page.path == "/download":

            # Handle file download requests
            params = parse_qs(page.query)

            class_name = params.get("class", [None])[0]
            filename = params.get("file", [None])[0]
            filename = unquote(filename)

            # Basic validation
            if not class_name or not filename:
                self.send_error(400, "Bad request")
                return

            # Prevent path traversal
            if ".." in class_name or ".." in filename or "/" in filename:
                self.send_error(403, "Forbidden")
                return

            filepath = os.path.join("uploads", f"class_{class_name}", filename)

            if not os.path.isfile(filepath):
                self.send_error(404, "File not found")
                return

            # Serve file as download
            self.send_response(200)
            self.send_header("Content-Type", "application/octet-stream")
            self.send_header(
                "Content-Disposition",
                f'attachment; filename="{filename}"'
            )
            self.send_header("Content-Length", os.path.getsize(filepath))
            self.end_headers()

            with open(filepath, "rb") as f:
                self.wfile.write(f.read())
            return

        elif page.path == '/vocab':
            params = parse_qs(page.query)
            filename = params.get("file", [None])[0]
            self.send_response(200)
            self.send_header("Content-type", "text/html")
            self.end_headers()
            with open("vocab.html", "rb") as f:
                html = f.read()
            html = html
            self.wfile.write(html)

        else:
            # Serve main page
            self.send_response(200)
            self.send_header("Content-type", "text/html")
            self.end_headers()
            with open("index.html", "rb") as f:
                html = f.read()
            html = html.replace(b"<!-- ITEMS -->", make_html_list().encode('utf-8'))
            self.wfile.write(html)
        
    def do_POST(self):
        # Parse the form data posted
        content_length = int(self.headers['Content-Length'])
        content_type = self.headers.get('Content-Type')
        form = cgi.FieldStorage(
            fp=self.rfile,
            headers=self.headers,
            environ={'REQUEST_METHOD': 'POST',
                     'CONTENT_TYPE': content_type}
        )
        action = form.getvalue("action")

        # Add Vocab Object
        if action == "add":
            vocab_value = form.getvalue("vocab")
            a = Vocab(vocab_value)
            classList.append(a)
            write_to_file()

            self.send_response(200)
            self.send_header("Content-type", "text/html")
            self.end_headers()
            with open("index.html", "rb") as f:
                html = f.read()
                html = html.replace(b"<!-- ITEMS -->", make_html_list().encode('utf-8'))
                self.wfile.write(html)

        # Delete Vocab Object
        if action == "delete":
            id_value = int(form.getvalue("id", [-1])[0])
            # Delete the files associated with the class
            class_name = classList[id_value].name
            class_dir = f"uploads{os.sep}class_{class_name}"
            if os.path.isdir(class_dir):
                for filename in os.listdir(class_dir):
                    file_path = os.path.join(class_dir, filename)
                    if os.path.isfile(file_path):
                        os.remove(file_path)
                os.rmdir(class_dir)
            classList.pop(id_value)
            write_to_file()
            # Send user back to main page
            self.send_response(200)
            self.send_header("Content-type", "text/html")
            self.end_headers()
            with open("index.html", "rb") as f:
                html = f.read()
                html = html.replace(b"<!-- ITEMS -->", make_html_list().encode('utf-8'))
                self.wfile.write(html)
        
        # Delete Individual File
        if action == "delete_file":
            id_value = int(form.getvalue("id", [-1])[0])
            class_name = form.getvalue("classname", [-1])
            class_id = form.getvalue("classID")
            # Checks for error
            if not os.path.exists(f"uploads{os.sep}class_{class_name}"):
                return "Invalid class"
            file_array = os.listdir(f"uploads{os.sep}class_{class_name}")
            # Checks for valid ID value inside the file array
            if 0 <= id_value < len(file_array):
                file_name = file_array[id_value] #Retrieves the File Name
                file_dir = f"uploads{os.sep}class_{class_name}"
                # Delete the individual files associated with the class
                os.remove(os.path.join(file_dir, file_name))
                #Send user back to Class page and Reload the page
                self.send_response(200)
                self.send_header("Content-type", "text/html")
                self.end_headers()
                with open("class.html", "rb") as f:
                    html = f.read()
                    html = html.replace(b"<!-- FILES -->", make_file_html_list(class_name, class_id).encode('utf-8'))
                    html = html.replace(b"CLASSIDPLACEHOLDER", class_id.encode('utf-8'))
                    self.wfile.write(html)


        # Storing and parsing file uploads

        if action == "upload":
            #write file to uploads directory within the class specified
            class_name = form.getvalue("class_name")
            class_number = form.getvalue("classID")
            os.makedirs(f"uploads{os.sep}class_{class_name}", exist_ok=True)
            file_data = form['file']
            filename = file_data.filename
            if file_data.filename:
                with open(f"uploads{os.sep}class_{class_name}{os.sep}{filename}", "wb") as f:
                    f.write(file_data.file.read())
            ext = Path(file_data.filename).suffix.lower()
            text = ""
            match ext:
                case ".pdf":
                    text = pypdf.PdfReader(f"uploads{os.sep}class_{class_name}{os.sep}{file_data.filename}")
                    text = "".join(page.extract_text() for page in text.pages)
                case ".docx":
                    readtext = docx.Document(f"uploads{os.sep}class_{class_name}{os.sep}{file_data.filename}")
                    text = "".join(para.text for para in readtext.paragraphs)
                case ".doc":
                    print("DOC parsing not implemented yet.")
                case ".pptx":
                    readtext = Presentation(f"uploads{os.sep}class_{class_name}{os.sep}{file_data.filename}")
                    text = " ".join(paragraph.text for slide in readtext.slides 
                                for shape in slide.shapes
                                   if shape.has_text_frame
                                    for paragraph in shape.text_frame.paragraphs)
                    #print("PPTX parsing not implemented yet.")
                case ".ppt":
                    print("PPT parsing not implemented yet.")
                case ".txt":
                    text = open(f"uploads{os.sep}class_{class_name}{os.sep}{file_data.filename}").read()
                    # print("TXT parsing not implemented yet.")
            # Sanitize the text
            text = text.lower()
            text = ''.join([i for i in text if i.isalpha() or i.isspace()])

            #GREAT MERGER sanitize and remove existing words from text
            text = classList[int(class_number)].sanitize_text(text)

            first_word = text.split(" ")[0]
            rest_of_text = " ".join(text.split()[1:])

            self.send_response(200)
            self.send_header("Content-type", "text/html")
            self.end_headers()
            with open("adder.html", "rb") as f:
                html = f.read()
                html = html.replace(b"CLASSIDPLACEHOLDER", class_number.encode('utf-8'))
                html = html.replace(b"FILENAMEPLACEHOLDER", filename.encode('utf-8'))
                html = html.replace(b"RESTOFTEXTPLACEHOLDER", rest_of_text.encode('utf-8'))
                html = html.replace(b"WORDINQUESTION", first_word.encode('utf-8'))
                self.wfile.write(html)
        
        #Handle adding a word to a class vocab Object
        if action == "addword":
            # Grab basic data from the form
            class_number = form.getvalue("classID")
            filename = form.getvalue("filename")
            rest_of_text = form.getvalue("restOfText")
            root = form.getvalue("root")
            pos = form.getvalue("pos")
            include_in_vocab_list = form.getvalue("addToList")
            english = form.getvalue("englishMeaning")
            transitive = form.getvalue("transitive")
            plural = form.getvalue("plural")

            #make and add Root object
            rootObject = Root(root, pos, filename, include_in_vocab_list=include_in_vocab_list, english=english, plural=plural, transitive=transitive)
            classList[int(class_number)].add_root(rootObject)

            #removes the root and its forms from the proceeding text
            if (rest_of_text):
                rest_of_text = classList[int(class_number)].fast_sanitize(text = rest_of_text, root=rootObject)

            #grab the first word
            first_word = rest_of_text.split()[0]
            #reformats text without first word
            rest_of_text = rest_of_text[len(first_word):].strip()
            #saves the vocab data 
            write_to_file()

            self.send_response(200)
            self.send_header("Content-type", "text/html")
            self.end_headers()
            with open("adder.html", "rb") as f:
                html = f.read()
                # Hold classID in hidden value
                html = html.replace(b"CLASSIDPLACEHOLDER", class_number.encode('utf-8'))
                # Hold filename in hidden value
                html = html.replace(b"FILENAMEPLACEHOLDER", filename.encode('utf-8'))
                # Hold rest of text in hidden value
                html = html.replace(b"RESTOFTEXTPLACEHOLDER", rest_of_text.encode('utf-8'))
                # Set word in question to first word of rest of text
                html = html.replace(b"WORDINQUESTION", first_word.encode('utf-8'))
                self.wfile.write(html)
        

if __name__ == "__main__":
    server = HTTPServer(('localhost', port), SimpleHandler)
    print(f"Server running on http://localhost:{port}")
    server.serve_forever()
    #server.shutdown() #stops server in the program

